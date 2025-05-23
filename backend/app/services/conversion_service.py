import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import pdfplumber
from docx import Document
from PIL import Image
from backend.app.utils.file_utils import generate_unique_filename

class FileConversionService:
    @staticmethod
    def txt_to_pptx(txt_path: str, pptx_path: str):
        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        with open(txt_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        chunk_size = 10
        for i in range(0, len(lines), chunk_size):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            content = ''.join(lines[i:i+chunk_size])
            left = Inches(0.5)
            top = Inches(0.5)
            width = slide_width - Inches(1)
            height = slide_height - Inches(1)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = content
            p.font.size = Pt(20)
        prs.save(pptx_path)
        return pptx_path

    @staticmethod
    def pdf_to_pptx(pdf_path: str, pptx_path: str):
        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
                text = page.extract_text() or ''
                left = Inches(0.5)
                top = Inches(0.5)
                width = slide_width - Inches(1)
                height = slide_height - Inches(1)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = text
                p.font.size = Pt(20)
        prs.save(pptx_path)
        return pptx_path

    @staticmethod
    def docx_to_pptx(docx_path: str, pptx_path: str):
        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        doc = Document(docx_path)
        chunk = []
        chunk_size = 10
        for para in doc.paragraphs:
            if para.text.strip():
                chunk.append(para.text.strip())
                if len(chunk) >= chunk_size:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = slide_width - Inches(1)
                    height = slide_height - Inches(1)
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    tf = textbox.text_frame
                    tf.word_wrap = True
                    p = tf.add_paragraph()
                    p.text = '\n'.join(chunk)
                    p.font.size = Pt(20)
                    chunk = []
        if chunk:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            left = Inches(0.5)
            top = Inches(0.5)
            width = slide_width - Inches(1)
            height = slide_height - Inches(1)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = '\n'.join(chunk)
            p.font.size = Pt(20)
        prs.save(pptx_path)
        return pptx_path

    @staticmethod
    def images_to_pptx(image_paths: list, pptx_path: str):
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # blank
        for img_path in image_paths:
            slide = prs.slides.add_slide(blank_slide_layout)
            with Image.open(img_path) as img:
                width_px, height_px = img.size
            # Convert pixels to inches (assuming 96 DPI)
            width_inches = width_px / 96
            height_inches = height_px / 96
            width = Inches(width_inches)
            height = Inches(height_inches)
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            # Scale image to fit slide while maintaining aspect ratio
            scale = min(slide_width / width, slide_height / height, 1)
            img_width = width * scale
            img_height = height * scale
            left = (slide_width - img_width) / 2
            top = (slide_height - img_height) / 2
            slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)
            # Add caption (filename) below the image
            caption_text = os.path.basename(img_path)
            caption_top = top + img_height + Inches(0.2)
            caption_height = Inches(0.6)
            # Only add caption if there's space below the image
            if caption_top + caption_height < slide_height:
                caption_left = Inches(0.5)
                caption_width = slide_width - Inches(1)
                textbox = slide.shapes.add_textbox(caption_left, caption_top, caption_width, caption_height)
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = caption_text
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.italic = True
        prs.save(pptx_path)
        return pptx_path

    @staticmethod
    def convert_to_pptx(src_path: str, ext: str, dest_dir: str, batch_image_paths: list = None) -> str:
        # Generate PPTX filename based on original filename
        original_name = Path(src_path).stem
        pptx_name = generate_unique_filename(f"{original_name}.pptx")
        pptx_path = os.path.join(dest_dir, pptx_name)
        
        if ext.lower() == '.txt':
            return FileConversionService.txt_to_pptx(src_path, pptx_path)
        elif ext.lower() == '.pdf':
            return FileConversionService.pdf_to_pptx(src_path, pptx_path)
        elif ext.lower() in ['.doc', '.docx']:
            return FileConversionService.docx_to_pptx(src_path, pptx_path)
        elif ext.lower() in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
            # If batch, convert all images to one PPTX
            if batch_image_paths:
                return FileConversionService.images_to_pptx(batch_image_paths, pptx_path)
            else:
                return FileConversionService.images_to_pptx([src_path], pptx_path)
        else:
            raise ValueError(f'Unsupported file type for conversion: {ext}') 