from pptx import Presentation
from typing import List, Dict
import logging
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import os
import tempfile

logger = logging.getLogger("pptx_service")

class PPTXService:
    @staticmethod
    def extract_text_from_pptx(pptx_path: str) -> List[Dict[str, str]]:
        """
        Extract text content from a PPTX file, including slide numbers and content.
        Returns a list of dictionaries containing slide number and content.
        """
        try:
            prs = Presentation(pptx_path)
            slides_content = []
            
            for idx, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes (textboxes, etc.)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # If slide has any text, add it to the content
                if slide_text:
                    slides_content.append({
                        "slide_number": idx,
                        "content": "\n".join(slide_text)
                    })
            
            logger.info(f"Successfully extracted text from {len(slides_content)} slides in {pptx_path}")
            return slides_content
            
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {pptx_path}: {str(e)}", exc_info=True)
            raise

    @staticmethod
    def format_slides_for_prompt(slides_content: List[Dict[str, str]]) -> str:
        """
        Format slide content into a prompt-friendly string.
        """
        formatted_content = "Slide Deck Content:\n\n"
        for slide in slides_content:
            formatted_content += f"Slide {slide['slide_number']}:\n{slide['content']}\n\n"
        return formatted_content.strip()

    @staticmethod
    def extract_images_from_pptx(pptx_path: str) -> List[Dict[str, str]]:
        """
        Extract the first image from each slide in a PPTX file.
        Returns a list of dicts: {slide_number, image_path}
        """
        try:
            prs = Presentation(pptx_path)
            slides_images = []
            for idx, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        image_bytes = image.blob
                        ext = image.ext
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{ext}') as tmp:
                            tmp.write(image_bytes)
                            tmp_path = tmp.name
                        slides_images.append({
                            "slide_number": idx,
                            "image_path": tmp_path
                        })
                        break  # Only first image per slide for now
            logger.info(f"Extracted images from {len(slides_images)} slides in {pptx_path}")
            return slides_images
        except Exception as e:
            logger.error(f"Error extracting images from PPTX {pptx_path}: {str(e)}", exc_info=True)
            raise 

    @staticmethod
    def extract_image_from_slide(pptx_path: str, slide_number: int):
        """
        Extract the first image from the given slide number in a PPTX file.
        Returns (image_bytes, mime_type) or (None, None) if not found.
        """
        try:
            prs = Presentation(pptx_path)
            if slide_number < 1 or slide_number > len(prs.slides):
                return None, None
            slide = prs.slides[slide_number - 1]
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext.lower()
                    if ext == 'jpeg' or ext == 'jpg':
                        mime_type = 'image/jpeg'
                    elif ext == 'png':
                        mime_type = 'image/png'
                    elif ext == 'bmp':
                        mime_type = 'image/bmp'
                    elif ext == 'gif':
                        mime_type = 'image/gif'
                    else:
                        mime_type = 'application/octet-stream'
                    return image_bytes, mime_type
            return None, None
        except Exception as e:
            logger.error(f"Error extracting image from slide {slide_number} in PPTX {pptx_path}: {str(e)}", exc_info=True)
            return None, None 